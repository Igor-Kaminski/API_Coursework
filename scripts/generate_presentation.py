"""Generate the coursework presentation PowerPoint."""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

OUTPUT = Path(__file__).resolve().parent.parent / "docs" / "presentation.pptx"

BG = RGBColor(0x0F, 0x17, 0x2A)
ACCENT = RGBColor(0x38, 0xBD, 0xF8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xBB, 0xBB, 0xBB)
DARK_CARD = RGBColor(0x1A, 0x25, 0x3C)
GREEN = RGBColor(0x4A, 0xDE, 0x80)
ORANGE = RGBColor(0xFB, 0xBF, 0x24)
RED = RGBColor(0xF8, 0x71, 0x71)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, left, top, width, height, text, size=18, color=WHITE,
             bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                     Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return tf


def add_bullet_list(slide, left, top, width, height, items, size=16,
                    color=WHITE, bullet_color=ACCENT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                     Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(6)
        p.space_before = Pt(2)

        bullet_run = p.add_run()
        bullet_run.text = "  \u25B8  "
        bullet_run.font.size = Pt(size)
        bullet_run.font.color.rgb = bullet_color
        bullet_run.font.bold = True

        text_run = p.add_run()
        text_run.text = item
        text_run.font.size = Pt(size)
        text_run.font.color.rgb = color
    return tf


def add_card(slide, left, top, width, height, title, body_lines,
             title_color=ACCENT):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_CARD
    shape.line.fill.background()
    shape.shadow.inherit = False

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(12)
    tf.margin_right = Pt(12)
    tf.margin_top = Pt(8)
    tf.margin_bottom = Pt(8)

    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(15)
    p.font.color.rgb = title_color
    p.font.bold = True
    p.space_after = Pt(6)

    for line in body_lines:
        p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(12)
        p.font.color.rgb = LIGHT_GRAY
        p.space_before = Pt(2)


def add_section_title(slide, text):
    add_text(slide, 0.6, 0.3, 8.8, 0.6, text, size=28, color=ACCENT, bold=True)
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.6), Inches(0.85), Inches(1.5), Pt(3),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()


def build():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    # ── Slide 1: Title ──
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(sl, BG)
    add_text(sl, 0.6, 1.0, 8.8, 1.0,
             "Rail Reliability & Delay Analytics API",
             size=36, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(sl, 0.6, 2.0, 8.8, 0.5,
             "COMP3011 Web Services and Web Data \u2014 Coursework 1",
             size=18, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
    add_text(sl, 0.6, 2.8, 8.8, 0.4,
             "Igor Kaminski",
             size=20, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(sl, 0.6, 3.5, 8.8, 0.4,
             "github.com/Igor-Kaminski/API_Coursework",
             size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
    add_text(sl, 0.6, 3.9, 8.8, 0.4,
             "Live: rail-api-coursework.onrender.com",
             size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    # ── Slide 2: Project Overview ──
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(sl, BG)
    add_section_title(sl, "Project Overview")
    add_text(sl, 0.6, 1.2, 8.8, 0.6,
             "A RESTful API for analysing UK rail service reliability and delays,\n"
             "powered by real operational data from the National Rail Darwin feed.",
             size=16, color=LIGHT_GRAY)
    add_card(sl, 0.6, 2.2, 4.2, 1.5, "Read-Only Analytics", [
        "Route reliability & average delay",
        "Cancellation rates & delay distributions",
        "Station hotspots & delay reasons",
        "Top delayed / cancelled routes",
    ])
    add_card(sl, 5.2, 2.2, 4.2, 1.5, "CRUD Incidents", [
        "User-reported service disruptions",
        "Role-based create / update / delete",
        "Filtered by route, station, severity",
        "Independent of imported data",
    ])
    add_text(sl, 0.6, 4.0, 8.8, 0.5,
             "Dataset: ~2,993 stations  \u2022  ~2,899 routes  \u2022  ~18,588 journey records",
             size=14, color=ACCENT, alignment=PP_ALIGN.CENTER)

    # ── Slide 3: Tech Stack ──
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(sl, BG)
    add_section_title(sl, "Technology Stack")

    techs = [
        ("FastAPI", "Modern async Python framework\nAuto OpenAPI / Swagger docs"),
        ("PostgreSQL", "Relational DB with FK constraints\nNatural fit for the data model"),
        ("SQLAlchemy 2.0", "Type-safe ORM\nMaintainable query layer"),
        ("Alembic", "Schema migrations\nVersioned & reproducible"),
        ("Pydantic v2", "Request/response validation\nExplicit schemas"),
        ("pytest", "81% line coverage\nIsolated DB fixtures"),
    ]
    for i, (name, desc) in enumerate(techs):
        col = i % 3
        row = i // 3
        add_card(sl, 0.6 + col * 3.1, 1.2 + row * 1.6, 2.9, 1.4, name,
                 desc.split("\n"))

    # ── Slide 4: Architecture ──
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(sl, BG)
    add_section_title(sl, "Architecture")
    add_text(sl, 0.6, 1.2, 8.8, 0.4,
             "Modular monolith \u2014 simple to test, explain, and deploy",
             size=16, color=LIGHT_GRAY)

    layers = [
        ("Routers", "stations \u2022 routes \u2022 incidents \u2022 analytics"),
        ("Schemas", "Pydantic request / response models"),
        ("Services", "Analytics engine \u2022 Import pipelines"),
        ("Models", "SQLAlchemy ORM \u2192 PostgreSQL"),
    ]
    for i, (name, desc) in enumerate(layers):
        y = 1.9 + i * 0.75
        shape = sl.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(1.5), Inches(y), Inches(7.0), Inches(0.6),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = DARK_CARD
        shape.line.fill.background()
        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Pt(10)
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = f"{name}   "
        r.font.size = Pt(14)
        r.font.color.rgb = ACCENT
        r.font.bold = True
        r2 = p.add_run()
        r2.text = desc
        r2.font.size = Pt(13)
        r2.font.color.rgb = LIGHT_GRAY

        if i < len(layers) - 1:
            arrow = sl.shapes.add_shape(
                MSO_SHAPE.DOWN_ARROW,
                Inches(4.8), Inches(y + 0.58), Inches(0.4), Inches(0.18),
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = ACCENT
            arrow.line.fill.background()

    add_text(sl, 0.6, 4.6, 4.0, 0.4,
             "Also: MCP server over same DB & logic",
             size=14, color=GREEN)
    add_text(sl, 5.0, 4.6, 4.5, 0.4,
             "core/  config \u2022 database \u2022 security \u2022 errors",
             size=13, color=LIGHT_GRAY)

    # ── Slide 5: Data Model ──
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(sl, BG)
    add_section_title(sl, "Data Model")

    entities = [
        ("Station", "name, code, CRS, TIPLOC\ncity, lat/lng", 0.4),
        ("Route", "origin \u2192 destination\noperator, distance_km", 3.4),
        ("JourneyRecord", "scheduled vs actual times\nstatus, delay_minutes", 0.4),
        ("Incident", "title, description, type\nseverity, status", 3.4),
    ]
    colors = [ACCENT, GREEN, ORANGE, RED]
    for i, (name, desc, x) in enumerate(entities):
        row = i // 2
        add_card(sl, x, 1.2 + row * 1.8, 2.8, 1.5, name,
                 desc.split("\n"), title_color=colors[i])

    add_text(sl, 6.5, 1.2, 3.2, 3.5,
             "Key design decisions:\n\n"
             "\u25B8 Reference data kept separate\n   from operational history\n\n"
             "\u25B8 Incidents are user-generated,\n   never auto-derived from delays\n\n"
             "\u25B8 FK constraints enforce\n   referential integrity\n\n"
             "\u25B8 CRS + TIPLOC codes for\n   real-world station lookup",
             size=13, color=LIGHT_GRAY)

    # ── Slide 6: Data Ingestion ──
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(sl, BG)
    add_section_title(sl, "Data Ingestion Strategy")
    add_text(sl, 0.6, 1.2, 8.8, 0.4,
             "External feeds are data sources, not runtime dependencies",
             size=16, color=LIGHT_GRAY, bold=True)

    add_card(sl, 0.6, 1.9, 4.2, 1.3, "Darwin (Push Port)", [
        "Real-time snapshot feed via FTP",
        "Journey records with timing data",
        "Main source for delay analytics",
    ], title_color=ORANGE)
    add_card(sl, 5.2, 1.9, 4.2, 1.3, "KnowledgeBase", [
        "Station reference enrichment",
        "CRS codes, cities, coordinates",
        "Human-readable route names",
    ], title_color=GREEN)

    add_text(sl, 0.6, 3.5, 8.8, 0.4,
             "Challenges solved during ingestion:",
             size=15, color=WHITE, bold=True)
    add_bullet_list(sl, 0.6, 3.9, 8.8, 1.5, [
        "Code aliases & incomplete metadata \u2192 merge + deduplication logic",
        "Routes stored as internal codes \u2192 resolved to human-readable names",
        "1,959 / 2,899 routes fully resolved; 120 remain unresolved (junctions/depots)",
    ], size=13)

    # ── Slide 7: API Design ──
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(sl, BG)
    add_section_title(sl, "API & Authentication")
    add_text(sl, 0.6, 1.2, 5.0, 0.4,
             "Versioned under /api/v1  \u2022  4 endpoint groups",
             size=15, color=LIGHT_GRAY)

    add_card(sl, 0.6, 1.8, 4.8, 2.0, "Endpoint Groups", [
        "GET  /stations  (+ /code/{code})",
        "GET  /routes  (+ /code/{code})",
        "CRUD /incidents",
        "GET  /analytics/routes/{id}/reliability",
        "GET  /analytics/stations/hotspots",
        "GET  /analytics/delay-reasons/common",
        "... and 12 more analytics endpoints",
    ])

    add_card(sl, 5.8, 1.8, 3.8, 2.0, "Auth Model (X-API-Key)", [
        "admin   \u2192 full access",
        "operator \u2192 manage incidents",
        "user    \u2192 create incidents",
        "no key  \u2192 public reads",
        "",
        "Swagger Authorize button",
        "Error codes per endpoint",
    ], title_color=GREEN)

    add_text(sl, 0.6, 4.2, 8.8, 0.5,
             "Business identifiers: filter stations by code, name, city, CRS, TIPLOC \u2014 "
             "not just numeric IDs",
             size=13, color=LIGHT_GRAY)

    # ── Slide 8: Analytics ──
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(sl, BG)
    add_section_title(sl, "Analytics Endpoints")

    analytics = [
        ("Route Reliability", "on-time %, delayed %, cancelled %"),
        ("Average Delay", "mean delay minutes per route"),
        ("Cancellation Rate", "cancelled / total journeys"),
        ("Delay Distribution", "buckets: 0-5, 5-15, 15-30, 30+ min"),
        ("Station Hotspots", "most affected stations by delay"),
        ("Top Delayed Routes", "ranked by average delay"),
        ("Top Cancelled Routes", "ranked by cancellation rate"),
        ("Delay Reasons", "most common reasons (frequency)"),
        ("Incident Frequency", "incidents per date bucket"),
        ("Severity Breakdown", "low / medium / high / critical"),
    ]
    for i, (name, desc) in enumerate(analytics):
        col = i % 2
        row = i // 2
        x = 0.6 + col * 4.7
        y = 1.15 + row * 0.75
        tf = add_text(sl, x, y, 4.5, 0.35, "", size=12, color=WHITE)
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = f"{name}  "
        r.font.size = Pt(13)
        r.font.color.rgb = ACCENT
        r.font.bold = True
        r2 = p.add_run()
        r2.text = desc
        r2.font.size = Pt(12)
        r2.font.color.rgb = LIGHT_GRAY

    add_text(sl, 0.6, 4.7, 8.8, 0.5,
             "Design principle: analytics match the quality and timespan of available data",
             size=14, color=ORANGE, bold=True, alignment=PP_ALIGN.CENTER)

    # ── Slide 9: Error Handling & Security ──
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(sl, BG)
    add_section_title(sl, "Error Handling & Security")

    add_card(sl, 0.6, 1.2, 4.4, 2.2, "Standardised Errors", [
        "Shared ErrorResponse model",
        "401  Unauthorized",
        "403  Forbidden",
        "404  Not Found",
        "409  Conflict (FK violations)",
        "422  Validation Error",
        "500 / 503  Server errors",
    ])
    add_card(sl, 5.3, 1.2, 4.3, 2.2, "Security Measures", [
        "API key authentication (3 roles)",
        "Rate limiting (slowapi, 100/min)",
        "CORS middleware",
        "DB health checks at startup",
        "Structured exception handlers",
        "Conflict checks before delete",
        "",
    ], title_color=GREEN)

    add_text(sl, 0.6, 3.8, 8.8, 0.6,
             "Intentionally proportionate for coursework scope \u2014\n"
             "not a production security system, but not absent either.",
             size=14, color=LIGHT_GRAY)

    # ── Slide 10: Testing & Deployment ──
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(sl, BG)
    add_section_title(sl, "Testing & Deployment")

    add_card(sl, 0.6, 1.2, 4.4, 2.5, "Testing (pytest)", [
        "81% line coverage (pytest-cov)",
        "Isolated DB fixtures per test",
        "",
        "Covers:",
        "  Incident API permissions",
        "  Station & route lookups",
        "  Analytics endpoints",
        "  Import services & Darwin parsing",
        "  MCP tools & security",
    ])
    add_card(sl, 5.3, 1.2, 4.3, 2.5, "Deployment (Render)", [
        "render.yaml blueprint",
        "Web service + managed PostgreSQL",
        "Auto-generated API keys",
        "Alembic runs before startup",
        "",
        "Live at:",
        "rail-api-coursework.onrender.com",
        "",
        "Schema + data are separate concerns",
    ], title_color=GREEN)

    # ── Slide 11: Version Control & Deliverables ──
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(sl, BG)
    add_section_title(sl, "Version Control & Deliverables")

    add_card(sl, 0.6, 1.2, 4.4, 1.6, "Git & GitHub", [
        "Granular commit history on main",
        "Feature-by-feature progression",
        "Public repo with full visibility",
        "README with setup for Linux,",
        "  macOS, and Windows",
    ])
    add_card(sl, 5.3, 1.2, 4.3, 1.6, "Deliverables", [
        "Technical report (5-page PDF)",
        "API documentation (Swagger PDF)",
        "Presentation slides (this file)",
        "GenAI conversation logs",
        "  (ChatGPT + Cursor exports)",
    ], title_color=GREEN)

    add_card(sl, 0.6, 3.2, 9.0, 1.2, "MCP Extension (Advanced)", [
        "Same database & business logic exposed via Model Context Protocol",
        "stdio, SSE, and streamable-http transports  \u2022  role-based tool access",
        "Demonstrates the system through a second interface beyond HTTP",
    ], title_color=ORANGE)

    # ── Slide 12: Reflection ──
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(sl, BG)
    add_section_title(sl, "Reflection & Future Work")

    add_card(sl, 0.6, 1.2, 4.4, 1.6, "Challenges", [
        "Sourcing a suitable dataset",
        "Cleaning messy real-world data",
        "Route name resolution",
        "Scoping the project appropriately",
    ], title_color=ORANGE)
    add_card(sl, 5.3, 1.2, 4.3, 1.6, "Limitations", [
        "Short time window (~days, not months)",
        "Basic authentication model",
        "No long-term trend analysis",
        "Security not production-grade",
    ], title_color=RED)

    add_card(sl, 0.6, 3.2, 9.0, 1.2, "Future Improvements", [
        "Sample Darwin data over a wider date range for temporal trends",
        "Proper OAuth2 / JWT authentication for public deployment",
        "Frontend dashboard for visual analytics",
    ], title_color=GREEN)

    # ── Slide 13: Thank You ──
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(sl, BG)
    add_text(sl, 0.6, 1.5, 8.8, 0.8,
             "Thank you",
             size=36, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(sl, 0.6, 2.5, 8.8, 0.4,
             "Questions?",
             size=22, color=ACCENT, alignment=PP_ALIGN.CENTER)

    links = [
        "GitHub:  github.com/Igor-Kaminski/API_Coursework",
        "Live API:  rail-api-coursework.onrender.com",
        "Swagger:  rail-api-coursework.onrender.com/docs",
    ]
    for i, link in enumerate(links):
        add_text(sl, 0.6, 3.3 + i * 0.35, 8.8, 0.3,
                 link, size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    prs.save(str(OUTPUT))
    print(f"Done: {OUTPUT}")


if __name__ == "__main__":
    build()
